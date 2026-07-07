"""
World News 5分ニュース番組 .pptx 自動生成スクリプト
依存: python-pptx (インストール済み)
使用法: python make_news_pptx.py
出力:  world_news_20260516.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

BASE      = Path(__file__).parent
FILES_DIR = BASE / "files"
FILES_DIR.mkdir(exist_ok=True)
OUTPUT    = FILES_DIR / "world_news_20260516.pptx"

# ── スライドサイズ: 16:9 ──────────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

# ── 色 ────────────────────────────────────────────────────────────────────────
BG    = RGBColor(0x04, 0x04, 0x0C)
BG2   = RGBColor(0x09, 0x09, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN  = RGBColor(0x00, 0xC8, 0xFF)
YEL   = RGBColor(0xFF, 0xD7, 0x00)
RED   = RGBColor(0xFF, 0x30, 0x50)
ORG   = RGBColor(0xFF, 0x90, 0x40)
GRN   = RGBColor(0x30, 0xFF, 0x80)
DIM   = RGBColor(0x50, 0x50, 0x78)
RISK  = RGBColor(0xFF, 0x80, 0x80)

# ── スライドコンテンツ ────────────────────────────────────────────────────────
SLIDES = [
    dict(
        cover=True, color=CYAN,
        badge="OPENING",
        title="World News",
        subtitle="2026年5月16日　主要5ニュース",
        source_line="ソース: NHK・産経・日経・CNBC・FT・BBC・NYT・Reuters ほか",
        bullets=[], risk="",
    ),
    dict(
        cover=False, color=YEL,
        badge="① 米中首脳会談",
        title="「G2」の衝撃と台湾最強警告",
        subtitle="習主席「対応誤れば両国衝突・非常に危険」/ トランプ「約束せず」",
        source_line="NHK 5/14〜16・産経 5/14・CNBC・FT・BBC・NYT・Reuters",
        bullets=[
            "5/14〜15  北京で米中首脳会談（2017年以来初の訪中）",
            "習主席「対応誤れば両国衝突・非常に危険な状況に陥る」と最強けん制",
            "トランプ「約束せず」── 台湾への武器売却 140億ドルの行方が焦点",
            "産経「G2時代の予感」── 日本の安保への影響が懸念される",
        ],
        risk="▲ リスク：数週間以内に台湾武器売却の可否が判明。"
             "米中G2合意で日本が蚊帳の外になる恐れ。",
    ),
    dict(
        cover=False, color=RED,
        badge="② 米イラン戦争",
        title="交渉崩壊リスクとホルムズ封鎖",
        subtitle="トランプ「もう我慢しない」/ ブレント原油 $109 / 日本エネルギー危機",
        source_line="NHK 5/15・Reuters 5/15・Bloomberg 5/15・毎日 5/16",
        bullets=[
            "2月28日開戦から2.5ヶ月── 停戦交渉が行き詰まり（5/15現在）",
            "トランプ「もうこれ以上長く我慢しない」── 戦闘再開を示唆（5/15）",
            "オマーン沖でインド船籍が攻撃── ホルムズ海峡は依然不安定",
            "ブレント原油 $109（+2.18%）。日本「6月に代替調達7割確保」の見通し",
        ],
        risk="▲ リスク【最緊急】：数日〜数週間以内に戦闘再開の可能性。"
             "ホルムズ長期封鎖で日本エネルギー危機。",
    ),
    dict(
        cover=False, color=CYAN,
        badge="③ 日米関係",
        title="ベッセント来日とG2への懸念",
        subtitle="高市首相・電話会談で米中会談の内容説明受ける",
        source_line="NHK 5/11〜15・産経 5/14・The Diplomat 5/11",
        bullets=[
            "ベッセント財務長官来日（5/11〜12）── レアアース・イラン対応を協議",
            "米中会談後、日米首脳が電話会談（5/15夜）で内容を詳細共有",
            "産経「G2時代── トランプ台湾譲歩に懸念強く」",
            "レアアースETF（REMX） -3.6% 急落── 供給不安が継続",
        ],
        risk="▲ リスク：150日間関税暫定措置の期限（2026年夏）と"
             "台湾武器売却が日米関係の試金石。",
    ),
    dict(
        cover=False, color=ORG,
        badge="④ 台湾問題",
        title="習主席の最強警告",
        subtitle="「衝突すれば非常に危険」── トランプ「いかなる約束もしない」",
        source_line="NHK 5/14〜16・CNBC・FT・BBC・NYT・PBS・CBS・Reuters・The Guardian",
        bullets=[
            "習主席「対応誤れば衝突・非常に危険」── CNBC・FT・BBC・NYT 一斉報道",
            "トランプ「約束しなかった」── 台湾武器売却 $14B は「近く判断」",
            "高市「台湾有事」答弁後、日本閣僚が初めて中国を訪問（黄川田大臣）",
            "台湾は防衛費 $25B 可決済み── 独自抑止力の強化を継続",
        ],
        risk="▲ リスク：武器売却凍結→台湾抑止力低下→中国の圧力増大。"
             "数週間〜数ヶ月が分岐点。",
    ),
    dict(
        cover=False, color=GRN,
        badge="⑤ 自動車産業",
        title="ホンダ上場来初赤字・EV撤退加速",
        subtitle="4,239億円赤字 / 脱ガソリン撤回 / 日産EV中止 / トヨタはインドへ",
        source_line="産経 5/14・日経 5/14・NHK 5/11〜12",
        bullets=[
            "ホンダ26年3月期：上場来初の赤字 4,239億円。「脱ガソリン」撤回（産経 5/14）",
            "日産：米国EV 2車種の生産中止。ホンダも米国EV 3車種を開発・販売中止（NHK 5/12）",
            "日本郵船：ホルムズ封鎖で中東向け自動車輸送の代替ルートを検討（NHK 5/11）",
            "トヨタ：インドに新工場建設発表（2029年前半稼働）（NHK 5/11）",
        ],
        risk="▲ リスク：EV撤退の間にBYDが市場獲得。"
             "ホルムズ長期封鎖で輸送コストが恒常化する恐れ。",
    ),
    dict(
        cover=True, color=CYAN,
        badge="END",
        title="World News",
        subtitle="以上、本日の5大ニュースをお伝えしました",
        source_line="2026年5月16日",
        bullets=[], risk="",
    ),
]


# ── ユーティリティ ────────────────────────────────────────────────────────────
def rgb(r, g, b):
    return RGBColor(r, g, b)


def set_bg(slide, color: RGBColor):
    """スライド背景を単色塗りつぶしに設定"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=None):
    """塗りつぶし矩形を追加"""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p     = tf.paragraphs[0]
    p.alignment = align
    run   = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "メイリオ"
    if color:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = WHITE
    return txBox


def add_bullet_box(slide, bullets, left, top, width, height,
                   color=WHITE, font_size=16):
    """箇条書きテキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "● " + bullet
        run.font.size  = Pt(font_size)
        run.font.name  = "メイリオ"
        run.font.color.rgb = color
    return txBox


# ── スライド生成 ──────────────────────────────────────────────────────────────
def build_cover_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG)

    color = data["color"]
    M = Inches

    # ヘッダーバー
    add_rect(slide, M(0), M(0), SW, M(0.55), fill_color=BG2)
    add_rect(slide, M(0), M(0.55), SW, Pt(2), fill_color=color)
    add_text(slide, "[ World News ]",
             M(0.15), M(0.08), M(4), M(0.45),
             font_size=20, bold=True, color=WHITE)
    add_text(slide, "2026.05.16  LIVE",
             M(10.5), M(0.12), M(3), M(0.4),
             font_size=13, color=DIM)

    # バッジ
    bw = M(1.4)
    add_rect(slide, M(0.15), M(0.75), bw, M(0.32),
             fill_color=BG2, line_color=color, line_width=Pt(1))
    add_text(slide, data["badge"],
             M(0.2), M(0.78), bw, M(0.3),
             font_size=14, bold=True, color=color)

    # メインタイトル
    add_text(slide, data["title"],
             M(0), M(2.2), SW, M(1.5),
             font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # サブタイトル
    add_text(slide, data["subtitle"],
             M(0), M(3.6), SW, M(0.9),
             font_size=24, bold=False, color=color, align=PP_ALIGN.CENTER)

    # 区切り線
    add_rect(slide, M(4.0), M(4.5), M(5.33), Pt(1.5), fill_color=color)

    # フッター
    add_rect(slide, M(0), M(7.25), SW, M(0.25), fill_color=BG2)
    add_text(slide, data["source_line"],
             M(0.15), M(7.28), M(12), M(0.22),
             font_size=9, color=DIM)


def build_news_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG)

    color = data["color"]
    M = Inches

    # ── ヘッダーバー ──
    add_rect(slide, M(0), M(0), SW, M(0.55), fill_color=BG2)
    add_rect(slide, M(0), M(0.55), SW, Pt(2), fill_color=color)
    add_text(slide, "[ World News ]",
             M(0.15), M(0.08), M(4), M(0.45),
             font_size=18, bold=True, color=WHITE)
    add_text(slide, "2026.05.16  LIVE",
             M(10.5), M(0.12), M(3), M(0.4),
             font_size=12, color=DIM)

    # ── バッジ ──
    bw = M(min(len(data["badge"]) * 0.18 + 0.5, 3.8))
    add_rect(slide, M(0.15), M(0.68), bw, M(0.33),
             fill_color=BG2, line_color=color, line_width=Pt(1))
    add_text(slide, data["badge"],
             M(0.22), M(0.71), bw, M(0.3),
             font_size=14, bold=True, color=color)

    # ── タイトル ──
    add_text(slide, data["title"],
             M(0.15), M(1.1), M(13), M(0.9),
             font_size=30, bold=True, color=WHITE)

    # ── サブタイトル ──
    add_text(slide, data["subtitle"],
             M(0.15), M(2.0), M(13), M(0.6),
             font_size=15, bold=False, color=color)

    # ── 区切り線 ──
    add_rect(slide, M(0.15), M(2.62), M(13), Pt(1), fill_color=color)

    # ── 箇条書き ──
    add_bullet_box(slide, data["bullets"],
                   M(0.15), M(2.75), M(13), M(2.8),
                   color=rgb(180, 180, 220), font_size=15)

    # ── リスクフッター ──
    if data.get("risk"):
        add_rect(slide, M(0), M(6.5), SW, M(1.0), fill_color=rgb(18, 4, 8))
        add_rect(slide, M(0), M(6.5), SW, Pt(1.5), fill_color=RED)
        add_text(slide, data["risk"],
                 M(0.15), M(6.58), M(13), M(0.85),
                 font_size=13, bold=False, color=RISK)

    # ── ソースフッター ──
    add_rect(slide, M(0), M(7.25), SW, M(0.25), fill_color=BG2)
    add_text(slide, data["source_line"],
             M(0.15), M(7.28), M(13), M(0.22),
             font_size=9, color=DIM)


# ── メイン ────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("=" * 60)
    print("World News PowerPoint 生成を開始します")
    print(f"出力先: {OUTPUT}")
    print("=" * 60)

    for i, data in enumerate(SLIDES):
        label = data["badge"]
        print(f"  [{i+1}/{len(SLIDES)}] スライド作成: {label}")
        if data["cover"]:
            build_cover_slide(prs, data)
        else:
            build_news_slide(prs, data)

    prs.save(str(OUTPUT))
    print(f"\n[完了] {OUTPUT}")
    print(f"   スライド数: {len(SLIDES)} 枚")


if __name__ == "__main__":
    main()
