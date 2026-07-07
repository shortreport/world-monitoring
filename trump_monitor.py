#!/usr/bin/env python3
"""
trump_monitor.py
================
トランプ大統領・主要閣僚の言動モニタリングツール

当日と昨日のトランプ大統領および主要閣僚の言動を
複数のソースから収集し、Claude AI で日本語要約します。

データソース:
  1. Truth Social          : トランプ大統領のSNS公式投稿
  2. ホワイトハウス公式RSS  : プレスリリース・声明・大統領令
  3. YouTube (WH公式)      : 記者会見・ブリーフィング動画（字幕付き）
  4. Google News RSS       : メディア報道・インタビュー記事
  5. World News API        : 詳細ニュース本文

使い方:
  python trump_monitor.py            # 通常実行
  python trump_monitor.py --save     # レポートをファイルに保存
  python trump_monitor.py --verbose  # 収集内容の詳細表示

必要パッケージ（venv に導入済）:
  anthropic, requests, beautifulsoup4, feedparser, rich, youtube-transcript-api
"""

import argparse
import json
import logging
import os
import re
import sys
import time
import urllib.parse
from datetime import datetime, date, timezone, timedelta
from pathlib import Path
from typing import Optional

# Windows UTF-8 対応
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

# ── ネットワーク待機（スリープ解除直後のDNS失敗対策）────────────────────────
import socket as _socket
def _wait_for_network(host="api.anthropic.com", port=443, timeout=60):
    import time as _t
    deadline = _t.time() + timeout
    while _t.time() < deadline:
        try:
            _socket.setdefaulttimeout(5)
            _socket.getaddrinfo(host, port)
            return True
        except OSError:
            _t.sleep(5)
    return False

if not _wait_for_network():
    sys.exit("[ERROR] ネットワーク接続タイムアウト（60秒）")

# ── サードパーティ ────────────────────────────────────────────────────────────
try:
    import anthropic
    import feedparser
    import requests
    from bs4 import BeautifulSoup
    from rich import box
    from rich.console import Console
    from rich.panel import Panel
    from rich.progress import Progress, SpinnerColumn, TextColumn
    from rich.table import Table
except ImportError as e:
    print(f"❌ 必要なパッケージが不足しています: {e}")
    print("   venv\\Scripts\\pip install anthropic requests beautifulsoup4 feedparser rich")
    sys.exit(1)

try:
    from youtube_transcript_api import YouTubeTranscriptApi
    TRANSCRIPT_OK = True
except ImportError:
    TRANSCRIPT_OK = False

# ── API キー ─────────────────────────────────────────────────────────────────
try:
    from api_config import ANTHROPIC_API_KEY, WORLDNEWS_API_KEY
except ImportError:
    ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY", "")
    WORLDNEWS_API_KEY = os.getenv("WORLDNEWS_API_KEY", "")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 設定 / Configuration
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

TODAY      = datetime.now(timezone.utc).date()
YESTERDAY  = TODAY - timedelta(days=1)
DATE_FROM  = datetime.combine(YESTERDAY, datetime.min.time()).replace(tzinfo=timezone.utc)
DATE_TO    = datetime.now(timezone.utc)
DATE_RANGE = f"{YESTERDAY} ~ {TODAY}"

# Claude モデル（要約用）
MODEL = "claude-haiku-4-5-20251001"

# ホワイトハウス公式 YouTube チャンネル ID（@WhiteHouse）
WH_YT_CHANNEL = "UCYxRlFDqcWM4y7FfpiAN3KQ"

# 監視対象：主要閣僚・要人
TARGETS = [
    "JD Vance",        # 副大統領
    "Marco Rubio",     # 国務長官
    "Pete Hegseth",    # 国防長官
    "Scott Bessent",   # 財務長官
    "Pam Bondi",       # 司法長官
    "Kristi Noem",     # 国土安全保障長官
    "Stephen Miller",  # 上席政策顧問
    "Russell Vought",  # OMB 局長
]

# ホワイトハウス RSS フィード
WH_RSS_FEEDS = {
    "briefings":           "https://www.whitehouse.gov/briefings-statements/feed/",
    "news":                "https://www.whitehouse.gov/news/feed/",
    "presidential_actions": "https://www.whitehouse.gov/presidential-actions/feed/",
}

# Google News 検索クエリ（RSS）
GNEWS_QUERIES = [
    "Trump press conference statement",
    "Trump interview media",
    "White House press briefing Karoline Leavitt",
    "Trump executive order announcement",
    '"JD Vance" OR "Marco Rubio" OR "Pete Hegseth" statement',
    '"Scott Bessent" OR "Pam Bondi" OR "Kristi Noem" statement',
]

console = Console(legacy_windows=False)
logger  = logging.getLogger(__name__)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ユーティリティ / Utilities
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def in_range(dt) -> bool:
    """datetime/date が対象期間内かどうかを返す"""
    d = dt.date() if isinstance(dt, datetime) else dt
    return d in (TODAY, YESTERDAY)

def strip_html(html: str) -> str:
    return BeautifulSoup(html, "html.parser").get_text(" ", strip=True)

def fetch_html(url: str, timeout: int = 12) -> Optional[str]:
    try:
        r = requests.get(
            url,
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"},
            timeout=timeout,
        )
        r.raise_for_status()
        return r.text
    except Exception as e:
        logger.debug("fetch_html %s: %s", url, e)
        return None

def extract_article_text(html: str, max_chars: int = 4000) -> str:
    """HTML から主要な本文テキストを抽出する"""
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup.find_all(["script", "style", "nav", "header", "footer", "aside", "form"]):
        tag.decompose()
    body = (
        soup.find("article")
        or soup.find("main")
        or soup.find("div", class_=re.compile(r"article|content|body|post|entry", re.I))
        or soup.find("body")
    )
    text = body.get_text("\n", strip=True) if body else ""
    return text[:max_chars]

def fmt_date(dt: datetime) -> str:
    return dt.strftime("%Y-%m-%d %H:%M UTC")

def parsed_to_utc(t) -> Optional[datetime]:
    """feedparser の time.struct_time → UTC datetime"""
    if t is None:
        return None
    return datetime(*t[:6], tzinfo=timezone.utc)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 1. Truth Social ─ トランプ大統領の公式 SNS 投稿
#    Truth Social API は匿名アクセスをブロックしているため、
#    ① Truth Social API（認証なし）→ ② 主要メディアによる Truth Social 報道
#    の順でフォールバックして取得する。
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _collect_truth_social_api(username: str = "realDonaldTrump") -> list[dict]:
    """Truth Social 公開 API からトランプ投稿を取得（403 の場合は空リストを返す）"""
    base    = "https://truthsocial.com/api/v1"
    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
    try:
        r = requests.get(f"{base}/accounts/lookup", params={"acct": username},
                         headers=headers, timeout=15)
        if r.status_code != 200:
            logger.debug("Truth Social API lookup HTTP %s", r.status_code)
            return []
        account_id = r.json().get("id")
        if not account_id:
            return []
        r2 = requests.get(f"{base}/accounts/{account_id}/statuses",
                          params={"limit": 40}, headers=headers, timeout=15)
        if r2.status_code != 200:
            return []
        posts = []
        for s in r2.json():
            raw = s.get("created_at", "")
            try:
                dt = datetime.fromisoformat(raw.replace("Z", "+00:00"))
            except ValueError:
                continue
            if not in_range(dt):
                continue
            text = strip_html(s.get("content", ""))
            if not text:
                continue
            posts.append({
                "source":  "Truth Social（直接取得）",
                "type":    "sns",
                "author":  "Donald Trump",
                "date":    fmt_date(dt),
                "title":   f"Truth Social投稿 ({dt.strftime('%m/%d %H:%M UTC')})",
                "content": text,
                "url":     s.get("url", ""),
            })
        return posts
    except Exception as e:
        logger.debug("Truth Social API error: %s", e)
        return []

def _collect_truth_social_via_news() -> list[dict]:
    """Truth Social の投稿を報道経由（Google News RSS）で収集するフォールバック"""
    queries = [
        "Trump Truth Social post today",
        "Trump posts Truth Social statement",
    ]
    items = []
    seen: set[str] = set()
    for q in queries:
        for art in _gnews_fetch(q, n=5):
            if art["url"] in seen:
                continue
            seen.add(art["url"])
            content = _fetch_article_text(art["url"])
            items.append({
                "source":  f"Truth Social報道 ({art['source']})",
                "type":    "sns",
                "author":  "Donald Trump（報道経由）",
                "date":    art["date"],
                "title":   art["title"],
                "content": content,
                "url":     art["url"],
            })
        time.sleep(0.3)
    return items

def collect_truth_social() -> list[dict]:
    """Truth Social 投稿を取得（直接→報道経由の順でフォールバック）"""
    posts = _collect_truth_social_api()
    if posts:
        return posts
    # API が403でブロックされている場合はニュース報道経由で収集
    logger.debug("Truth Social API blocked, falling back to news search")
    return _collect_truth_social_via_news()


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 2. ホワイトハウス公式 RSS ─ 声明・大統領令・ニュース
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

TYPE_LABEL = {
    "briefings":            "ブリーフィング・公式声明",
    "news":                 "ホワイトハウスニュース",
    "presidential_actions": "大統領令・行政措置",
}

def collect_whitehouse_rss() -> list[dict]:
    items = []
    for feed_key, feed_url in WH_RSS_FEEDS.items():
        try:
            feed = feedparser.parse(feed_url)
            for e in feed.entries:
                dt = parsed_to_utc(e.get("published_parsed"))
                if dt is None or not in_range(dt):
                    continue

                # 全文取得（whitehouse.gov はスクレイプ可能）
                link    = e.get("link", "")
                content = ""
                if link:
                    html = fetch_html(link)
                    if html:
                        content = extract_article_text(html, 5000)
                if not content:
                    content = strip_html(e.get("summary", ""))

                items.append({
                    "source":  "ホワイトハウス公式",
                    "type":    feed_key,
                    "author":  "White House",
                    "date":    fmt_date(dt),
                    "title":   e.get("title", "無題"),
                    "content": content,
                    "url":     link,
                })
        except Exception as ex:
            logger.warning("WH RSS [%s] error: %s", feed_key, ex)
    return items


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 3. YouTube ─ ホワイトハウス公式チャンネルの動画 + 字幕
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def get_yt_transcript(video_id: str) -> str:
    if not TRANSCRIPT_OK:
        return ""
    try:
        segments = YouTubeTranscriptApi.get_transcript(video_id, languages=["en", "en-US"])
        return " ".join(s["text"] for s in segments)[:10000]
    except Exception as e:
        logger.debug("transcript %s: %s", video_id, e)
        return ""

def collect_whitehouse_youtube() -> list[dict]:
    rss_url = f"https://www.youtube.com/feeds/videos.xml?channel_id={WH_YT_CHANNEL}"
    try:
        feed = feedparser.parse(
            rss_url,
            request_headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"},
        )
        videos = []
        for e in feed.entries:
            dt = parsed_to_utc(e.get("published_parsed"))
            if dt is None or not in_range(dt):
                continue

            # video ID の取得
            vid_id = e.get("yt_videoid", "")
            if not vid_id:
                m = re.search(r"v=([^&]+)", e.get("link", ""))
                vid_id = m.group(1) if m else ""

            transcript = get_yt_transcript(vid_id) if vid_id else ""
            content    = transcript or strip_html(e.get("summary", ""))

            videos.append({
                "source":         "YouTube (ホワイトハウス公式)",
                "type":           "video",
                "author":         "White House",
                "date":           fmt_date(dt),
                "title":          e.get("title", "無題"),
                "content":        content,
                "url":            e.get("link", f"https://www.youtube.com/watch?v={vid_id}"),
                "has_transcript": bool(transcript),
            })
        return videos
    except Exception as ex:
        logger.warning("YouTube RSS error: %s", ex)
        return []


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 4. Google News RSS ─ メディア報道・インタビュー・記者会見カバレッジ
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _gnews_fetch(query: str, n: int = 6) -> list[dict]:
    q   = urllib.parse.quote(query)
    url = f"https://news.google.com/rss/search?q={q}&hl=en&gl=US&ceid=US:en"
    try:
        feed = feedparser.parse(url, request_headers={"User-Agent": "TrumpMonitor/1.0"})
        results = []
        for e in feed.entries[:n]:
            dt = parsed_to_utc(e.get("published_parsed"))
            if dt and not in_range(dt):
                continue
            results.append({
                "title":  e.get("title", "").strip(),
                "url":    e.get("link", ""),
                "source": e.get("source", {}).get("title", ""),
                "date":   fmt_date(dt) if dt else "",
            })
        return results
    except Exception as ex:
        logger.debug("gnews [%s]: %s", query, ex)
        return []

def _fetch_article_text(url: str) -> str:
    """記事URLからテキストを取得（失敗時は空文字）"""
    # Google News のリダイレクト URL は直接フェッチできないため除外
    if "news.google.com" in url:
        return ""
    html = fetch_html(url)
    return extract_article_text(html, 3000) if html else ""

def collect_google_news() -> list[dict]:
    seen_urls: set[str] = set()
    items: list[dict]   = []

    for query in GNEWS_QUERIES:
        for art in _gnews_fetch(query):
            url = art["url"]
            if url in seen_urls:
                continue
            seen_urls.add(url)
            content = _fetch_article_text(url)
            items.append({
                "source":  art["source"] or "News",
                "type":    "news",
                "author":  "",
                "date":    art["date"],
                "title":   art["title"],
                "content": content,
                "url":     url,
            })
        time.sleep(0.4)  # rate limiting

    return items


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 5. World News API ─ 詳細ニュース（本文付き）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

WNA_QUERIES = [
    "Trump press conference",
    "Trump interview",
    "White House press briefing",
    "Trump executive order",
    "JD Vance OR Marco Rubio OR Pete Hegseth statement",
    "Scott Bessent OR Pam Bondi statement",
]

_wna_warned = False  # 401 の重複警告を防ぐフラグ

def _wna_search(query: str, n: int = 5) -> list[dict]:
    global _wna_warned
    if not WORLDNEWS_API_KEY:
        return []
    try:
        r = requests.get(
            "https://api.worldnewsapi.com/search-news",
            params={
                "api-key":               WORLDNEWS_API_KEY,
                "text":                  query,
                "language":              "en",
                "number":                n,
                "earliest-publish-date": f"{YESTERDAY} 00:00:00",
                "sort":                  "publish-time",
                "sort-direction":        "DESC",
            },
            timeout=15,
        )
        if r.status_code == 401:
            if not _wna_warned:
                logger.warning("World News API: 認証エラー(401)。api_config.py の WORLDNEWS_API_KEY を確認してください。")
                _wna_warned = True
            return []
        if r.status_code != 200:
            logger.debug("WNA HTTP %s for [%s]", r.status_code, query)
            return []
        articles = []
        for a in r.json().get("news", []):
            raw = a.get("publish_date", "")
            try:
                dt      = datetime.fromisoformat(raw.replace("Z", "+00:00"))
                date_s  = fmt_date(dt)
                in_rng  = in_range(dt)
            except ValueError:
                date_s = raw
                in_rng = True  # 日付が不明な場合は含める
            if not in_rng:
                continue
            articles.append({
                "source":  a.get("url", "").split("/")[2] if a.get("url") else "WNA",
                "type":    "news",
                "author":  a.get("author", ""),
                "date":    date_s,
                "title":   a.get("title", "無題"),
                "content": (a.get("text") or a.get("summary") or "")[:3000],
                "url":     a.get("url", ""),
            })
        return articles
    except Exception as ex:
        logger.warning("WNA error [%s]: %s", query, ex)
        return []

def collect_world_news() -> list[dict]:
    seen_urls: set[str] = set()
    items: list[dict]   = []
    for q in WNA_QUERIES:
        for art in _wna_search(q):
            if art["url"] not in seen_urls:
                seen_urls.add(art["url"])
                items.append(art)
        time.sleep(0.3)
    return items


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 6. Claude AI 要約
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

SYSTEM = """\
あなたはアメリカ政治の専門アナリストです。
収集されたデータをもとにトランプ大統領と主要閣僚の言動を日本語で包括的に要約してください。

要約の原則:
- 重要な発言・決定を優先的に取り上げる
- 動画コンテンツ（字幕あり）は特に詳しく要約する
- 発言の背景・文脈・政策的意味を簡潔に添える
- 情報源の種類（SNS/記者会見/インタビュー/公式発表/動画）を明記する
- 日本の安全保障・経済・外交に関連する事項はハイライトする
- 事実を客観的に伝える（政治的評価・批判は避ける）\
"""

def _build_input(items: list[dict]) -> str:
    order = ["sns", "presidential_actions", "briefings", "video", "news"]
    buckets: dict[str, list[dict]] = {k: [] for k in order}
    for it in items:
        t = it.get("type", "news")
        buckets.setdefault(t, []).append(it)

    labels = {
        "sns":                  "📱 SNS投稿（Truth Social）",
        "presidential_actions": "📜 大統領令・行政措置",
        "briefings":            "🎤 ブリーフィング・公式声明",
        "video":                "🎬 動画コンテンツ（字幕付き優先）",
        "news":                 "📰 ニュース・インタビュー報道",
    }

    lines = []
    for key in order:
        bucket = buckets.get(key, [])
        if not bucket:
            continue
        lines.append(f"\n{'='*60}\n{labels.get(key, key)}\n{'='*60}")
        for i, it in enumerate(bucket, 1):
            lines.append(f"\n[{i}] {it.get('title', '無題')}")
            lines.append(f"    日時    : {it.get('date', '不明')}")
            lines.append(f"    情報源  : {it.get('source', '不明')}")
            if it.get("author"):
                lines.append(f"    著者    : {it['author']}")
            if it.get("has_transcript"):
                lines.append("    ★字幕あり（動画全文）")
            lines.append(f"    URL     : {it.get('url', 'N/A')}")
            body = it.get("content", "").strip()
            if body:
                lines.append(f"    内容:\n{body[:2500]}")
    return "\n".join(lines)

def summarize(items: list[dict]) -> str:
    if not items:
        return "収集できたデータがありませんでした。ネットワーク接続または API キーを確認してください。"

    client  = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    data    = _build_input(items)

    prompt = f"""\
以下は {YESTERDAY}〜{TODAY}（UTC）に収集したトランプ大統領と主要閣僚に関する情報です（計 {len(items)} 件）。

{data}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
上記データをもとに、以下の形式で日本語レポートを作成してください。
「情報なし」のセクションは「（該当情報なし）」と記載してください。

## ■ 本日の主要ポイント（箇条書き 3〜6 項目）

## 1. トランプ大統領のSNS発言
主要な投稿を引用しながら要約。

## 2. トランプ大統領の記者会見・インタビュー
記者会見・メディアインタビューの発言内容を具体的に要約。

## 3. ホワイトハウス公式発表・大統領令
政策発表、大統領令、行政措置の内容と意義。

## 4. ホワイトハウス記者会見（報道官ブリーフィング）
Karoline Leavitt 等の報道官による記者会見の要点。

## 5. 動画コンテンツ（★字幕ありの動画を優先して詳述）
YouTube 動画のトランスクリプトに基づく詳細な内容要約。

## 6. 主要閣僚の動向
閣僚別（JD Vance、Marco Rubio、Pete Hegseth 等）の発言・動向をまとめる。

## 7. 日本への影響・注目点
日本の安全保障・経済・貿易・外交に直結する事項をハイライト。\
"""

    resp = client.messages.create(
        model=MODEL,
        max_tokens=6000,
        system=SYSTEM,
        messages=[{"role": "user", "content": prompt}],
    )
    return resp.content[0].text


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 6b. 構造化 JSON 要約（URL 付き）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

SYSTEM_STRUCTURED = """\
あなたはアメリカ政治の専門アナリストです。
収集されたデータをもとにトランプ大統領と主要閣僚の言動を日本語で分析し、
必ず有効な JSON のみを出力してください（コードブロック・説明文は不要）。

各 bullet は {"text": "日本語テキスト", "source": "情報源名", "url": "元データのURL"} 形式。
URL は入力データ内の「URL     :」フィールドから最も関連性の高いものを使用してください。
URLが不明または該当なしの場合は "" (空文字) にしてください。\
"""

def summarize_structured(items: list[dict]) -> dict:
    """Claude に URL 付き構造化 JSON で要約させる（generate_site.py 連携用）"""
    if not items:
        return {}

    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    data   = _build_input(items)

    prompt = f"""\
以下は {YESTERDAY}〜{TODAY}（UTC）に収集したトランプ大統領と主要閣僚に関する情報です（計 {len(items)} 件）。
各アイテムには「URL     :」フィールドがあります。各 bullet の url はそのフィールドから使用してください。

{data}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
上記データを分析し、以下の JSON スキーマのみを出力してください（コードブロック不要）。
情報がないセクションは空配列 [] にしてください。

{{
  "key_points":          [{{"text":"...","source":"...","url":"..."}}],
  "sns":                 [{{"text":"...","source":"...","url":"..."}}],
  "press_conference":    [{{"text":"...","source":"...","url":"..."}}],
  "whitehouse_official": [{{"text":"...","source":"...","url":"..."}}],
  "wh_briefing":         [{{"text":"...","source":"...","url":"..."}}],
  "video":               [{{"text":"...","source":"...","url":"..."}}],
  "cabinet":             [{{"text":"...","source":"...","url":"..."}}],
  "japan_impact":        [{{"text":"...","source":"...","url":"..."}}]
}}

件数の目安:
  key_points 3〜6件 / sns 最大6件 / press_conference 最大5件
  whitehouse_official 最大5件 / wh_briefing 最大4件 / video 最大4件
  cabinet 最大5件（閣僚名を冒頭に明記） / japan_impact 最大5件（【カテゴリ】を冒頭に付与）
"""

    resp = client.messages.create(
        model=MODEL,
        max_tokens=8000,
        system=SYSTEM_STRUCTURED,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = resp.content[0].text.strip()
    # コードブロックが付いていた場合は除去
    raw = re.sub(r'^```(?:json)?\s*\n?', '', raw)
    raw = re.sub(r'\n?```\s*$',          '', raw)

    try:
        return json.loads(raw)
    except json.JSONDecodeError as ex:
        logger.warning("summarize_structured JSON parse error: %s\nRaw(500): %s", ex, raw[:500])
        # テキスト中の最初の { ... } ブロックを抽出して再試行
        m = re.search(r'\{.*\}', raw, re.DOTALL)
        if m:
            try:
                return json.loads(m.group())
            except json.JSONDecodeError:
                pass
        console.print(f"[red]構造化JSON生成失敗: {ex}[/red]")
        return {}


def save_structured_json(structured: dict, items: list[dict]) -> Path:
    """URL 付き構造化 JSON を trump_latest.json として保存（generate_site.py が読み込む）"""
    PPTX_DIR.mkdir(parents=True, exist_ok=True)

    # ソース別件数を集計
    type_counts: dict[str, int] = {}
    for it in items:
        t = it.get("type", "news")
        type_counts[t] = type_counts.get(t, 0) + 1

    sources_summary = []
    for key, label, icon in [
        ("sns",                  "Truth Social",       "\U0001f4f1"),
        ("presidential_actions", "大統領令・行政措置", "\U0001f4dc"),
        ("briefings",            "ホワイトハウスRSS",          "\U0001f3a4"),
        ("video",                "YouTube（字幕あり）",             "\U0001f3ac"),
        ("news",                 "Google News",                                              "\U0001f4f0"),
    ]:
        cnt = type_counts.get(key, 0)
        if cnt > 0:
            sources_summary.append({"icon": icon, "label": label, "count": f"{cnt}件"})

    output = {
        "generated_at":    datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "date_range":      DATE_RANGE,
        "item_count":      len(items),
        "sources_summary": sources_summary,
        **structured,
    }

    out_path = PPTX_DIR / "trump_latest.json"
    out_path.write_text(json.dumps(output, ensure_ascii=False, indent=2), encoding="utf-8")
    return out_path


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 表示・保存 / Display & Save
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def print_summary_table(items: list[dict]) -> None:
    t = Table(box=box.SIMPLE_HEAVY, show_header=True, header_style="bold cyan")
    t.add_column("カテゴリ",         style="cyan",       width=22)
    t.add_column("件数",             style="bold white", width=5, justify="right")
    t.add_column("主なソース",        style="dim",        width=40)

    order = [
        ("sns",                  "SNS投稿"),
        ("presidential_actions", "大統領令・行政措置"),
        ("briefings",            "ブリーフィング"),
        ("video",                "動画コンテンツ"),
        ("news",                 "ニュース・インタビュー"),
    ]
    for key, label in order:
        bucket = [x for x in items if x.get("type") == key]
        if bucket:
            srcs = list({x["source"] for x in bucket})[:3]
            t.add_row(label, str(len(bucket)), "、".join(srcs))
    console.print(t)

def print_verbose(items: list[dict]) -> None:
    """--verbose 時に各記事のタイトル・URL を表示"""
    for it in items:
        console.print(
            f"  [{it['type']:20s}] {it['title'][:60]:60s}  {it['date'][:16]}",
            style="dim"
        )

def save_report(summary: str, items: list[dict]) -> str:
    ts       = datetime.now().strftime("%Y%m%d_%H%M")
    filename = f"trump_report_{ts}.txt"
    with open(filename, "w", encoding="utf-8") as f:
        f.write("トランプ大統領・閣僚動向レポート\n")
        f.write(f"生成日時 : {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        f.write(f"対象期間 : {YESTERDAY} ～ {TODAY}（UTC）\n")
        f.write(f"収集件数 : {len(items)} 件\n")
        f.write("=" * 70 + "\n\n")
        f.write(summary)
        f.write("\n\n" + "=" * 70 + "\n【情報源一覧】\n")
        for it in items:
            f.write(f"- [{it.get('type','?'):20s}] {it.get('title','無題')[:70]}\n")
            if it.get("url"):
                f.write(f"  {it['url']}\n")
    return filename


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PowerPoint 生成 / PPTX Generation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

try:
    from pptx import Presentation as _Prs
    from pptx.util import Inches as _In, Pt as _Pt
    from pptx.dml.color import RGBColor as _RGB
    from pptx.enum.text import PP_ALIGN as _PA
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

# ── スライドサイズ 16:9 ───────────────────────────────────────────────────────
_SW = _In(13.33) if PPTX_OK else None
_SH = _In(7.50)  if PPTX_OK else None

# ── カラーパレット（白背景テーマ）───────────────────────────────────────────
def _c(r, g, b): return _RGB(r, g, b)
_BG    = _c(0xFF, 0xFF, 0xFF)   # 白背景
_BG2   = _c(0x16, 0x16, 0x32)   # ヘッダーバー（濃紺）
_FOOT  = _c(0xF0, 0xF0, 0xF5)   # フッターバー（薄グレー）
_WHITE = _c(0xFF, 0xFF, 0xFF)   # ヘッダー内白文字
_DARK  = _c(0x14, 0x14, 0x32)   # タイトル文字（濃紺）
_RED   = _c(0xC8, 0x1E, 0x32)   # 赤（白背景用）
_BLUE  = _c(0x3C, 0x3B, 0x6E)   # 星条旗ブルー（ENDスライド用）
_LBLUE = _c(0x00, 0x96, 0xC8)   # ライトブルー（白背景用）
_GOLD  = _c(0xB4, 0x82, 0x00)   # ゴールド（白背景用）
_BODY  = _c(0x32, 0x32, 0x50)   # 本文テキスト（濃グレー紺）
_DIM   = _c(0x8C, 0x8C, 0xA0)   # ソース行・日付

# セクション別アクセントカラー（スライド番号順）
_SEC_COLORS = [_GOLD, _LBLUE, _RED, _GOLD, _LBLUE, _RED, _GOLD, _RED]


# ── 低レベルヘルパー ──────────────────────────────────────────────────────────

def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = lw
    else:
        s.line.fill.background()
    return s

def _text(slide, text, l, t, w, h, size=16, bold=False,
          color=None, align=None):
    if align is None:
        align = _PA.LEFT
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = str(text)
    run.font.size      = _Pt(size)
    run.font.bold      = bold
    run.font.name      = "メイリオ"
    run.font.color.rgb = color or _WHITE
    return tb

def _bullets(slide, items, l, t, w, h, color=None, size=14):
    color = color or _WHITE
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment    = _PA.LEFT
        p.space_before = _Pt(10)
        run = p.add_run()
        run.text           = "▸  " + str(b)
        run.font.size      = _Pt(size)
        run.font.name      = "メイリオ"
        run.font.color.rgb = color
    return tb


# ── Markdown パーサー ─────────────────────────────────────────────────────────

def _clean(text: str) -> str:
    """Markdown 記号を除去してプレーンテキストにする"""
    text = re.sub(r'\*{1,3}(.+?)\*{1,3}', r'\1', text, flags=re.DOTALL)
    text = re.sub(r'\[(.+?)\]\(.+?\)',     r'\1', text)
    text = re.sub(r'`(.+?)`',              r'\1', text)
    text = re.sub(r'^#{1,6}\s+',           '',    text, flags=re.MULTILINE)
    text = re.sub(r'^\|.+',               '',    text, flags=re.MULTILINE)
    text = re.sub(r'\s{2,}',              ' ',   text)
    return text.strip()

def _trunc(text: str, n: int = 70) -> str:
    """タイトル等の短い文字列用（スペースが限られる場所専用）"""
    return text[:n] + "…" if len(text) > n else text

def _trunc_sent(text: str, max_len: int = 200) -> str:
    """箇条書き用：文末（。！？.!?）で区切り、文章が途切れないようにする"""
    if len(text) <= max_len:
        return text
    # max_len 以内で最後の文末文字を探す
    for ch in ('。', '！', '？', '.', '!', '?'):
        pos = text.rfind(ch, 0, max_len)
        if pos > max_len // 3:          # 先頭付近すぎる切断は避ける
            return text[:pos + 1]
    # 文末が見つからない場合は読点・句点・スペースで区切る
    for ch in ('、', '，', ' ', '　'):
        pos = text.rfind(ch, 0, max_len)
        if pos > max_len // 2:
            return text[:pos] + "…"
    return text[:max_len] + "…"

def _get_bullets(lines: list[str], max_n: int = 6) -> list[str]:
    """セクションの行リストから箇条書きを抽出（最大 max_n 件）"""
    bullets: list[str] = []
    for line in lines:
        s = line.strip()
        if s.startswith(("- ", "* ")):
            t = _clean(s[2:])
            if t and len(t) > 3:
                bullets.append(_trunc_sent(t))
        elif re.match(r'^\d+\.\s', s):
            t = _clean(re.sub(r'^\d+\.\s+', '', s))
            if t:
                bullets.append(_trunc_sent(t))
    # 箇条書きゼロ → 段落から抽出
    if not bullets:
        for line in lines:
            s = line.strip()
            if (s and not s.startswith(('#', '|', '---', '==='))
                    and len(s) > 10):
                t = _clean(s)
                if t:
                    bullets.append(_trunc_sent(t))
                if len(bullets) >= max_n:
                    break
    return bullets[:max_n]

def _parse_md(summary: str) -> list[dict]:
    """## ヘッダーでサマリーをセクションに分割"""
    sections: list[dict] = []
    cur = None
    for line in summary.split('\n'):
        if line.startswith('## '):
            if cur:
                sections.append(cur)
            cur = {'title': _clean(line), 'lines': []}
        elif cur is not None:
            cur['lines'].append(line)
    if cur:
        sections.append(cur)
    return sections


# ── スライドビルダー ──────────────────────────────────────────────────────────

def _slide_cover(prs, date_label: str, count: int):
    M = _In
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, _BG)
    # ヘッダーバー（濃紺）
    _rect(sl, M(0), M(0),    _SW, M(0.6),  fill=_BG2)
    _rect(sl, M(0), M(0.6),  _SW, _Pt(3),  fill=_RED)
    _text(sl, "TRUMP MONITOR",
          M(0.15), M(0.08), M(7), M(0.5),  size=22, bold=True, color=_WHITE)
    _text(sl, date_label,
          M(7.3), M(0.13),  M(5.9), M(0.42), size=13, color=_c(0xA0, 0xA0, 0xBE))
    # メインタイトル（白背景・濃紺文字）
    _text(sl, "Trump Monitor",
          M(0), M(1.8), _SW, M(1.6),
          size=64, bold=True, color=_DARK, align=_PA.CENTER)
    _text(sl, "トランプ大統領・主要閣僚 動向レポート",
          M(0), M(3.45), _SW, M(0.85),
          size=26, color=_GOLD, align=_PA.CENTER)
    _text(sl, date_label,
          M(0), M(4.35), _SW, M(0.55),
          size=18, color=_DIM, align=_PA.CENTER)
    _text(sl, f"収集件数 {count} 件  ／  AI分析: Claude {MODEL}",
          M(0), M(4.95), _SW, M(0.5),
          size=15, color=_DIM, align=_PA.CENTER)
    # フッター（薄グレー）
    _rect(sl, M(0), M(7.2),  _SW, M(0.3),  fill=_FOOT)
    _text(sl, "Source: White House RSS / Truth Social / YouTube / Google News",
          M(0.15), M(7.23), M(13), M(0.27), size=11, color=_DIM)

def _slide_content(prs, title: str, badge: str, bullets: list[str],
                   color, subtitle: str = ""):
    M = _In
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, _BG)
    # ヘッダーバー（濃紺）
    _rect(sl, M(0), M(0),    _SW, M(0.6),  fill=_BG2)
    _rect(sl, M(0), M(0.6),  _SW, _Pt(3),  fill=color)
    _text(sl, "TRUMP MONITOR",
          M(0.15), M(0.08), M(6), M(0.5),  size=20, bold=True, color=_WHITE)
    _text(sl, TODAY.strftime("%Y.%m.%d"),
          M(11.0), M(0.12), M(2.2), M(0.44), size=13, color=_c(0xA0, 0xA0, 0xBE))
    # バッジ（薄彩色背景）
    bw = M(min(len(badge) * 0.26 + 0.5, 5.5))
    ci = int(str(color), 16)
    cr, cg, cb = (ci >> 16) & 0xFF, (ci >> 8) & 0xFF, ci & 0xFF
    badge_fill = _c(
        min(255, 255 - (255 - cr) // 8),
        min(255, 255 - (255 - cg) // 8),
        min(255, 255 - (255 - cb) // 8),
    )
    _rect(sl, M(0.15), M(0.73), bw, M(0.38),
          fill=badge_fill, line=color, lw=_Pt(1.5))
    _text(sl, badge, M(0.23), M(0.76), bw, M(0.34),
          size=16, bold=True, color=color)
    # タイトル（濃紺文字）
    _text(sl, title, M(0.15), M(1.2), M(13), M(1.2),
          size=34, bold=True, color=_DARK)
    # 区切り線・箇条書き開始位置
    line_y, bullet_y = M(2.45), M(2.62)
    # 区切り線
    _rect(sl, M(0.15), line_y, M(13), _Pt(1.5), fill=color)
    # 箇条書き（最大5件、フォント21pt・濃グレー紺）
    _bullets(sl, bullets, M(0.15), bullet_y, M(13), M(4.3), color=_BODY, size=21)
    # フッター（薄グレー）
    _rect(sl, M(0), M(7.2),  _SW, M(0.3),  fill=_FOOT)
    _text(sl, "Source: White House / Truth Social / YouTube / Google News",
          M(0.15), M(7.23), M(13), M(0.27), size=11, color=_DIM)

def _slide_end(prs, date_label: str):
    M = _In
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, _BG)
    # ヘッダーバー（濃紺）
    _rect(sl, M(0), M(0),    _SW, M(0.6),  fill=_BG2)
    _rect(sl, M(0), M(0.6),  _SW, _Pt(3),  fill=_RED)
    _text(sl, "TRUMP MONITOR",
          M(0.15), M(0.08), M(6), M(0.5),  size=22, bold=True, color=_WHITE)
    # END テキスト（星条旗ブルー・白背景でも視認性あり）
    _text(sl, "END",
          M(0), M(2.5), _SW, M(1.7),
          size=90, bold=True, color=_BLUE, align=_PA.CENTER)
    _text(sl, "以上、本日のレポートをお届けしました",
          M(0), M(4.2), _SW, M(0.75),
          size=26, color=_GOLD, align=_PA.CENTER)
    _text(sl, date_label,
          M(0), M(5.0), _SW, M(0.55),
          size=16, color=_DIM, align=_PA.CENTER)
    # フッター（薄グレー）
    _rect(sl, M(0), M(7.2),  _SW, M(0.3),  fill=_FOOT)
    _text(sl, "Generated by Claude AI  ／  Source: WH RSS / Truth Social / YouTube / Google News",
          M(0.15), M(7.23), M(13), M(0.27), size=11, color=_DIM)


# ── メイン PPTX 生成関数 ────────────────────────────────────────────────────

# セクション番号 → (バッジラベル, アクセントカラー)
_BADGE_MAP = [
    ("■ 主要ポイント",      _GOLD),
    ("SNS 発言",            _LBLUE),
    ("記者会見・インタビュー", _RED),
    ("公式発表・大統領令",   _GOLD),
    ("WH 記者会見",         _LBLUE),
    ("動画コンテンツ",       _RED),
    ("閣僚動向",            _GOLD),
    ("日本への影響",         _RED),
]

# PPTX 保存先ディレクトリ
PPTX_DIR = Path("db") / "north america"

def generate_pptx(summary: str, items: list[dict]) -> Path:
    """Claude サマリーと収集データから PPTX を生成して PPTX_DIR に保存する"""
    if not PPTX_OK:
        raise ImportError("python-pptx が必要です: pip install python-pptx")

    PPTX_DIR.mkdir(parents=True, exist_ok=True)

    ts       = datetime.now().strftime("%Y%m%d_%H%M")
    out_path = PPTX_DIR / f"trump_report_{ts}.pptx"

    prs = _Prs()
    prs.slide_width  = _SW
    prs.slide_height = _SH

    date_label = f"対象期間: {YESTERDAY} ～ {TODAY}（UTC）"

    # 1. カバースライド
    _slide_cover(prs, date_label, len(items))

    # 2. コンテンツスライド（## セクション 1 枚ずつ）
    sections = _parse_md(summary)
    for i, sec in enumerate(sections):
        title   = sec['title']
        bullets = _get_bullets(sec['lines'], max_n=5)
        if not bullets:
            bullets = ["（該当情報なし）"]
        badge, color = _BADGE_MAP[i] if i < len(_BADGE_MAP) else (f"#{i+1}", _GOLD)
        _slide_content(prs, title, badge, bullets, color)

    # 3. エンドスライド
    _slide_end(prs, date_label)

    prs.save(str(out_path))
    return out_path


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# メイン / Main
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def main() -> None:
    parser = argparse.ArgumentParser(description="トランプ大統領・閣僚動向モニター")
    parser.add_argument("--save",    action="store_true", help="レポートをファイルに保存")
    parser.add_argument("--verbose", action="store_true", help="収集記事の詳細を表示")
    parser.add_argument("--debug",   action="store_true", help="デバッグログを有効化")
    args = parser.parse_args()

    logging.basicConfig(
        level=logging.DEBUG if args.debug else logging.WARNING,
        format="%(levelname)s [%(name)s]: %(message)s",
    )

    console.print(Panel(
        f"[bold cyan]トランプ大統領・閣僚動向モニター[/bold cyan]\n"
        f"対象期間: [yellow]{YESTERDAY}[/yellow] ～ [yellow]{TODAY}[/yellow]（UTC）\n"
        f"モデル  : [green]{MODEL}[/green]",
        box=box.DOUBLE, expand=False,
    ))
    console.print()

    all_items: list[dict] = []

    tasks_def = [
        ("📱 Truth Social 取得中",           collect_truth_social),
        ("🏛️  ホワイトハウス RSS 取得中",    collect_whitehouse_rss),
        ("🎬 YouTube 動画・字幕 取得中",      collect_whitehouse_youtube),
        ("📰 Google News 検索中",            collect_google_news),
        ("🌐 World News API 検索中",         collect_world_news),
    ]

    with Progress(SpinnerColumn(), TextColumn("{task.description}"), console=console) as prog:
        for label, fn in tasks_def:
            tid = prog.add_task(label, total=None)
            result = fn()
            all_items.extend(result)
            prog.update(tid, description=f"{label.split()[0]} {label.split()[1]} … [green]{len(result)}件[/green]")
            prog.stop_task(tid)

    console.print()
    print_summary_table(all_items)
    if args.verbose:
        console.print("\n[dim]── 収集記事一覧 ──[/dim]")
        print_verbose(all_items)
    console.print(f"\n合計 [bold green]{len(all_items)} 件[/bold green] を収集しました。")

    if not all_items:
        console.print("[red]⚠ データが収集できませんでした。ネットワーク接続を確認してください。[/red]")
        return

    console.print("\n[bold]🤖 Claude AI で要約中（数秒かかります）…[/bold]")
    report = summarize(all_items)

    console.print()
    console.print(Panel(report, title="[bold cyan]【要約レポート】[/bold cyan]",
                        box=box.ROUNDED, padding=(1, 2)))

    # テキストレポートを保存（--save オプション）
    if args.save:
        fname = save_report(report, all_items)
        console.print(f"\n[OK] テキストレポート: [green]{fname}[/green]")

    # PPTX を常に生成
    console.print("\n[bold]PowerPoint を生成中...[/bold]")
    try:
        pptx_path = generate_pptx(report, all_items)
        console.print(f"[OK] PowerPoint  : [green]{pptx_path}[/green]")
    except Exception as e:
        console.print(f"[red]PPTX 生成エラー: {e}[/red]")

    # URL 付き構造化 JSON を生成（generate_site.py 連携用）
    console.print("\n[bold]構造化 JSON を生成中（generate_site.py 連携用）...[/bold]")
    try:
        structured = summarize_structured(all_items)
        if structured:
            json_path = save_structured_json(structured, all_items)
            console.print(f"[OK] JSON        : [green]{json_path}[/green]")
            console.print(
                "     -> python generate_site.py を実行すると"
                " trump.html が実データで更新されます。"
            )
        else:
            console.print("[yellow]JSON 生成スキップ（構造化データが空）[/yellow]")
    except Exception as ex:
        console.print(f"[red]JSON 生成エラー: {ex}[/red]")


if __name__ == "__main__":
    main()
